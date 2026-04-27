"""
KindRide Investor Pitch Deck Generator
Generates a professional PowerPoint presentation for investor meetings.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x0D, 0x94, 0x88)   # primary brand
DARK_NAVY   = RGBColor(0x0F, 0x17, 0x2A)   # headlines
MID_SLATE   = RGBColor(0x33, 0x41, 0x55)   # body text
LIGHT_SLATE = RGBColor(0x94, 0xA3, 0xB8)   # captions
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFF)
ACCENT_RED  = RGBColor(0xDC, 0x26, 0x26)
LIGHT_TEAL  = RGBColor(0xCC, 0xFB, 0xF1)   # card background
SLIDE_BG    = RGBColor(0xF1, 0xF5, 0xF9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=SLIDE_BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), font_color=DARK_NAVY, bold=False,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox

def add_para(tf, text, font_size=Pt(14), font_color=MID_SLATE,
             bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p

def teal_bar(slide, height=Inches(0.08)):
    """Accent bar at the top of content slides."""
    add_rect(slide, 0, 0, SLIDE_W, height, TEAL)

def slide_header(slide, title, subtitle=None):
    teal_bar(slide)
    add_text(slide, title,
             Inches(0.6), Inches(0.22), Inches(12), Inches(0.7),
             font_size=Pt(32), font_color=DARK_NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.95), Inches(11), Inches(0.4),
                 font_size=Pt(15), font_color=TEAL, bold=False)

def bullet_card(slide, title, bullets, left, top, width, height, icon=""):
    add_rect(slide, left, top, width, height, WHITE,
             line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(1))
    # card title
    add_text(slide, f"{icon}  {title}" if icon else title,
             left + Inches(0.2), top + Inches(0.15),
             width - Inches(0.4), Inches(0.45),
             font_size=Pt(15), font_color=TEAL, bold=True)
    # divider
    add_rect(slide, left + Inches(0.2), top + Inches(0.62),
             width - Inches(0.4), Pt(1.5), TEAL)
    # bullets
    txBox = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.72),
        width - Inches(0.5), height - Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        add_para(tf, f"• {b}", font_size=Pt(12.5), font_color=MID_SLATE,
                 space_before=Pt(4 if i > 0 else 0))

def stat_box(slide, number, label, left, top, width=Inches(2.8), height=Inches(1.4)):
    add_rect(slide, left, top, width, height, TEAL)
    add_text(slide, number,
             left, top + Inches(0.15), width, Inches(0.65),
             font_size=Pt(34), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             left + Inches(0.1), top + Inches(0.82), width - Inches(0.2), Inches(0.5),
             font_size=Pt(12), font_color=WHITE, align=PP_ALIGN.CENTER)

def footer(slide, text="KindRide  ·  Confidential  ·  2025"):
    add_text(slide, text,
             Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             font_size=Pt(9), font_color=LIGHT_SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  TITLE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)

# large teal left accent bar
add_rect(s, 0, 0, Inches(0.5), SLIDE_H, TEAL)

# glow circle (decorative)
circ = s.shapes.add_shape(9, Inches(8.5), Inches(-1), Inches(6), Inches(6))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x0D, 0x94, 0x88)
circ.fill.fore_color.transparency = 0.88
circ.line.fill.background()

add_text(s, "KindRide",
         Inches(1.2), Inches(1.6), Inches(8), Inches(1.4),
         font_size=Pt(72), font_color=WHITE, bold=True)

add_text(s, "Ride-sharing built on community, trust, and zero detour.",
         Inches(1.2), Inches(3.05), Inches(9), Inches(0.7),
         font_size=Pt(22), font_color=TEAL)

add_rect(s, Inches(1.2), Inches(3.85), Inches(2.2), Pt(3), TEAL)

add_text(s, "Investor Briefing  ·  Series Seed",
         Inches(1.2), Inches(4.05), Inches(7), Inches(0.5),
         font_size=Pt(14), font_color=LIGHT_SLATE)

add_text(s, "Oluwafemi Adebayo  ·  Founder & CEO",
         Inches(1.2), Inches(4.6), Inches(7), Inches(0.4),
         font_size=Pt(13), font_color=LIGHT_SLATE)

add_text(s, "2025  ·  CONFIDENTIAL",
         Inches(1.2), Inches(6.9), Inches(5), Inches(0.35),
         font_size=Pt(10), font_color=RGBColor(0x47, 0x55, 0x69))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Problem", "Getting a ride shouldn't require a stranger willing to detour for you.")
footer(s)

# 3 problem cards
problems = [
    ("💸", "Prohibitively Expensive",
     ["Uber/Lyft surge pricing regularly hits 3–5×",
      "Low-income commuters often cannot afford reliable transport",
      "~45% of US adults report skipping trips due to cost"]),
    ("🌍", "Environmental Waste",
     ["Every ride-hail trip adds a dedicated vehicle to the road",
      "Avg. ride-hail deadhead rate: 40% of miles driven empty",
      "Transport = 29% of US greenhouse gas emissions"]),
    ("🤝", "Broken Trust in Sharing",
     ["Existing platforms commoditise the driver–passenger relationship",
      "Safety incidents erode confidence, especially for women & elderly",
      "No meaningful community layer — just transactions"]),
]
card_w = Inches(3.9)
for i, (icon, title, bullets) in enumerate(problems):
    bullet_card(s, title, bullets,
                Inches(0.55 + i * 4.18), Inches(1.55),
                card_w, Inches(5.0), icon)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  THE SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Solution — KindRide", "Connecting passengers with drivers already heading their way.")
footer(s)

# Hero statement box
add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.15), TEAL)
add_text(s,
         "KindRide matches passengers with drivers who are ALREADY travelling in the same direction — "
         "no detour, minimal cost, maximum trust.",
         Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.95),
         font_size=Pt(16), font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)

pillars = [
    ("Zero-Detour Model", [
        "Driver is already going your way — adds 0 miles",
        "Fare is a fraction of traditional ride-hail",
        "Carbon footprint: near zero marginal increase"]),
    ("Community Hubs", [
        "Campus, hospital, corporate & community networks",
        "Verified members build trust within groups",
        "Email-domain auto-join for institutions"]),
    ("Safety-First Platform", [
        "Live trip recording (privacy-preserving snapshots)",
        "One-tap SOS with emergency contact alerts",
        "ID verification via Stripe Identity"]),
    ("Kind Points Rewards", [
        "Drivers earn points redeemable for real value",
        "Reputation system rewards safe, kind behaviour",
        "Gamified tiers: Helper → Champion → Legend"]),
]
card_w = Inches(2.9)
for i, (title, bullets) in enumerate(pillars):
    bullet_card(s, title, bullets,
                Inches(0.45 + i * 3.1), Inches(2.82),
                card_w, Inches(4.2))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4  —  HOW IT WORKS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "How It Works", "Five steps — from request to reward.")
footer(s)

steps = [
    ("1", "Passenger Posts Need",
     "Opens app, shares destination direction. No scheduling, no waiting rooms."),
    ("2", "Smart Matching",
     "Algorithm surfaces drivers already headed that way, ranked by proximity, trust score & hub membership."),
    ("3", "Driver Accepts",
     "Driver reviews passenger reputation & accepts. Live GPS tracking begins for both parties."),
    ("4", "Trip Starts Automatically",
     "Proximity detection triggers trip start when driver arrives — no QR codes, no friction."),
    ("5", "Mutual Rating & Points",
     "Driver earns Kind Points. Both rate each other. Recordings auto-expire unless flagged."),
]

box_h = Inches(1.08)
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.5) + i * (box_h + Inches(0.08))
    # number bubble
    add_rect(s, Inches(0.5), y, Inches(0.6), box_h, TEAL)
    add_text(s, num, Inches(0.5), y + Inches(0.22), Inches(0.6), Inches(0.55),
             font_size=Pt(22), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # content box
    add_rect(s, Inches(1.2), y, Inches(11.5), box_h, WHITE,
             line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(1))
    add_text(s, title, Inches(1.4), y + Inches(0.07), Inches(4), Inches(0.42),
             font_size=Pt(14), font_color=DARK_NAVY, bold=True)
    add_text(s, desc, Inches(5.5), y + Inches(0.1), Inches(7), Inches(0.82),
             font_size=Pt(12.5), font_color=MID_SLATE)
    # connector arrow
    if i < len(steps) - 1:
        add_text(s, "▾", Inches(0.62), y + box_h, Inches(0.4), Inches(0.15),
                 font_size=Pt(9), font_color=TEAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5  —  MARKET OPPORTUNITY
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Market Opportunity", "A $285B global addressable market, underserved by incumbent platforms.")
footer(s)

# TAM / SAM / SOM nested boxes
add_rect(s, Inches(0.5),  Inches(1.5), Inches(6.5), Inches(5.5), RGBColor(0xCC, 0xFB, 0xF1))
add_rect(s, Inches(0.9),  Inches(1.9), Inches(5.7), Inches(4.6), RGBColor(0x99, 0xF6, 0xE4))
add_rect(s, Inches(1.3),  Inches(2.4), Inches(4.9), Inches(3.5), TEAL)

add_text(s, "TAM", Inches(0.6), Inches(1.55), Inches(2), Inches(0.4),
         font_size=Pt(11), font_color=RGBColor(0x0F, 0x76, 0x6E), bold=True)
add_text(s, "$285B  Global Ride-Hail Market (2024)",
         Inches(0.65), Inches(1.95), Inches(5.8), Inches(0.4),
         font_size=Pt(11), font_color=RGBColor(0x0F, 0x76, 0x6E))

add_text(s, "SAM", Inches(1.0), Inches(2.55), Inches(2), Inches(0.4),
         font_size=Pt(11), font_color=WHITE, bold=True)
add_text(s, "$42B  Community & commuter ride-sharing (US + UK)",
         Inches(1.05), Inches(2.98), Inches(4.7), Inches(0.4),
         font_size=Pt(11), font_color=WHITE)

add_text(s, "SOM", Inches(1.4), Inches(3.5), Inches(2), Inches(0.4),
         font_size=Pt(11), font_color=WHITE, bold=True)
add_text(s, "$1.2B  Year-5 realistic capture\n(campus + metro corridors)",
         Inches(1.45), Inches(3.92), Inches(4.3), Inches(0.65),
         font_size=Pt(11), font_color=WHITE)

# Right column stats
stats = [
    ("$8,000", "Average annual spend on\nprivate transport per US household"),
    ("680M",   "Global daily commuters\nunderserved by public transit"),
    ("23%",    "CAGR of ride-sharing market\n2024–2030 (Statista)"),
    ("4.2×",   "Higher retention in\ntrust-based platforms vs transactional"),
]
for i, (val, lbl) in enumerate(stats):
    r, c = divmod(i, 2)
    stat_box(s, val, lbl,
             Inches(7.4 + c * 3.0), Inches(1.55 + r * 1.65),
             Inches(2.7), Inches(1.45))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6  —  BUSINESS MODEL
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Business Model", "Multiple monetisation levers — none rely on exploiting drivers.")
footer(s)

streams = [
    ("Platform Fee\n(Primary)", "8–12% of suggested trip contribution from passenger.", "$$$"),
    ("Hub Subscriptions",       "Institutions (universities, hospitals) pay monthly SaaS fee for verified hub access.", "$$"),
    ("Kind Points Marketplace",  "Drivers redeem points via partner discounts; KindRide earns affiliate margin.", "$$"),
    ("Safety Premium",           "Optional monthly plan: extended recording retention, priority matching, badge.", "$"),
    ("Data Insights (B2B)",      "Anonymised mobility pattern reports sold to urban planners & transit authorities.", "$$"),
    ("Voluntary Tips",           "Stripe Connect tipping layer; KindRide takes 2.9% processing margin.", "$"),
]

card_w = Inches(3.9)
card_h = Inches(2.1)
for i, (title, desc, rev) in enumerate(streams):
    r, c = divmod(i, 3)
    left = Inches(0.5 + c * 4.2)
    top  = Inches(1.5 + r * 2.3)
    add_rect(s, left, top, card_w, card_h, WHITE,
             line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(1))
    add_rect(s, left, top, Inches(0.25), card_h, TEAL)
    add_text(s, title, left + Inches(0.35), top + Inches(0.12),
             card_w - Inches(0.9), Inches(0.52),
             font_size=Pt(13), font_color=DARK_NAVY, bold=True)
    add_text(s, desc, left + Inches(0.35), top + Inches(0.68),
             card_w - Inches(0.55), Inches(1.0),
             font_size=Pt(11.5), font_color=MID_SLATE)
    add_text(s, rev, left + card_w - Inches(0.65), top + Inches(0.1),
             Inches(0.5), Inches(0.4),
             font_size=Pt(15), font_color=TEAL, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7  —  TRACTION & PRODUCT STATUS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Traction & Product Status", "Full-stack product — live, tested, and iterated on real users.")
footer(s)

# Progress bars
milestones = [
    ("Core Ride Matching Engine",        1.0,  "Live"),
    ("Live GPS Tracking (Driver → Passenger)", 1.0, "Live"),
    ("Community Hubs + Domain Allow-list", 1.0, "Live"),
    ("Safety Recording (Privacy-Preserving)", 1.0, "Live"),
    ("SOS + Emergency Alerts",           1.0,  "Live"),
    ("Kind Points + Reputation System",  1.0,  "Live"),
    ("Trip Chat (Real-time)",            1.0,  "Live"),
    ("Identity Verification (Stripe)",   0.75, "Beta"),
    ("Stripe Tipping / Payments",        0.60, "In Progress"),
    ("Multi-leg Journey Handoffs",       0.85, "Beta"),
]

bar_w_total = Inches(5.5)
for i, (label, pct, status) in enumerate(milestones):
    y = Inches(1.55) + i * Inches(0.52)
    add_text(s, label, Inches(0.5), y, Inches(5.2), Inches(0.38),
             font_size=Pt(12), font_color=MID_SLATE)
    # track
    add_rect(s, Inches(5.8), y + Inches(0.08), bar_w_total, Inches(0.25),
             RGBColor(0xE2, 0xE8, 0xF0))
    # fill
    fill_color = TEAL if pct == 1.0 else RGBColor(0xF5, 0x9E, 0x0B)
    add_rect(s, Inches(5.8), y + Inches(0.08), Inches(bar_w_total.inches * pct), Inches(0.25), fill_color)
    # status badge
    badge_color = TEAL if status == "Live" else (RGBColor(0xF5, 0x9E, 0x0B) if status == "Beta" else RGBColor(0x64, 0x74, 0x8B))
    add_rect(s, Inches(11.55), y + Inches(0.03), Inches(1.3), Inches(0.32), badge_color)
    add_text(s, status, Inches(11.55), y + Inches(0.04), Inches(1.3), Inches(0.28),
             font_size=Pt(10), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8  —  TECHNOLOGY & SCALABILITY
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Technology & Scalability", "Cloud-native, serverless-ready stack built to handle 10× growth without re-architecting.")
footer(s)

tech_cols = [
    ("Mobile App", [
        "React Native (Expo) — single codebase, iOS + Android",
        "Expo Router file-based navigation",
        "Real-time Supabase subscriptions for GPS & chat",
        "Offline-resilient with local state fallback",
    ]),
    ("Backend API", [
        "FastAPI (Python) on Render — auto-scales horizontally",
        "Supabase PostgREST for low-latency data ops",
        "JWT-verified endpoints, Row Level Security on all tables",
        "Structured logging + ride lifecycle observability",
    ]),
    ("Infrastructure", [
        "Supabase (PostgreSQL + Realtime + Storage)",
        "Expo Push Notifications — no APNS/FCM dependency",
        "Stripe Identity + Connect for KYC & tipping",
        "Zero-trust: service role key never leaves backend",
    ]),
    ("Scale Path", [
        "Matching engine stateless → shardable by geography",
        "PostgIS geospatial queries already in schema",
        "Hub model enables viral city-by-city expansion",
        "Supabase edge functions for sub-50ms matching",
    ]),
]

card_w = Inches(3.0)
for i, (title, bullets) in enumerate(tech_cols):
    bullet_card(s, title, bullets,
                Inches(0.4 + i * 3.15), Inches(1.5),
                card_w, Inches(5.4))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9  —  COMPETITIVE LANDSCAPE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Competitive Landscape", "KindRide occupies a defensible white space no incumbent owns.")
footer(s)

headers = ["", "Uber / Lyft", "BlaBlaCar", "Waze Carpool", "KindRide"]
rows = [
    ("Zero-Detour Model",        "✗", "Partial", "Partial", "✓"),
    ("Community Hub Verified",   "✗", "✗",       "✗",       "✓"),
    ("Live GPS + Safety Rec.",    "✗", "✗",       "✗",       "✓"),
    ("Non-monetary Driver Reward","✗", "✗",       "✗",       "✓"),
    ("Multi-leg Journey Support", "✗", "✗",       "✗",       "✓"),
    ("Real-time Trip Chat",       "✗", "✗",       "✗",       "✓"),
    ("Avg. Passenger Cost",       "High", "Low", "Low",      "Very Low"),
    ("Driver Carbon Impact",      "High", "Low", "Low",      "Near Zero"),
]

col_w = Inches(2.3)
row_h = Inches(0.52)
table_left = Inches(0.4)
table_top  = Inches(1.5)

# header row
for j, h in enumerate(headers):
    bg = TEAL if j == 4 else DARK_NAVY
    add_rect(s, table_left + j * col_w, table_top, col_w, row_h, bg)
    add_text(s, h, table_left + j * col_w + Inches(0.05),
             table_top + Inches(0.1), col_w - Inches(0.1), Inches(0.38),
             font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    row_top = table_top + (i + 1) * row_h
    bg_row = OFF_WHITE if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        bg = LIGHT_TEAL if j == 4 else bg_row
        add_rect(s, table_left + j * col_w, row_top, col_w, row_h, bg,
                 line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(0.5))
        fc = TEAL if (j == 4 and cell == "✓") else (ACCENT_RED if cell == "✗" else MID_SLATE)
        fw = True if cell in ("✓", "✗") else False
        add_text(s, cell, table_left + j * col_w + Inches(0.05),
                 row_top + Inches(0.1), col_w - Inches(0.1), Inches(0.35),
                 font_size=Pt(12), font_color=fc, bold=fw, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10  —  CHALLENGES & MITIGATION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Challenges & How We Address Them",
             "We have identified every material risk and built mitigations into the product.")
footer(s)

challenges = [
    ("🐔", "Cold-Start / Chicken-and-Egg",
     "The classic two-sided marketplace problem.",
     ["Hub-first launch strategy seeds both sides simultaneously within closed communities",
      "Campus & corporate pilots provide captive supply AND demand from day one",
      "Founding Driver programme with permanent badge locks in early supply"]),
    ("⚖️", "Regulatory & Insurance",
     "Ride-sharing regulation varies by jurisdiction.",
     ["Zero-detour model positions KindRide as 'ridesharing' not 'for-hire transport' in most jurisdictions",
      "Driver insurance declaration captured at onboarding",
      "Legal review budgeted in seed round; advisory board includes transport law expertise"]),
    ("🔒", "Safety & Trust",
     "Single safety incident can destroy brand.",
     ["Privacy-preserving periodic recordings (not continuous) with 72h auto-expiry",
      "Stripe Identity KYC + community hub verification creates layered trust",
      "SOS one-tap emergency alert with SMS to named contacts"]),
    ("📈", "Retention & Habit Formation",
     "Users default to Uber when KindRide supply is thin.",
     ["Kind Points create genuine switching cost and habit loop",
      "Push notifications on driver match prevent app abandonment",
      "Multi-leg journey feature makes KindRide sticky for daily commuters"]),
]

for i, (icon, title, subtitle, mitigations) in enumerate(challenges):
    r, c = divmod(i, 2)
    left = Inches(0.4 + c * 6.45)
    top  = Inches(1.5 + r * 2.75)
    w, h = Inches(6.1), Inches(2.55)
    add_rect(s, left, top, w, h, WHITE,
             line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(1))
    add_rect(s, left, top, Inches(0.3), h, TEAL)
    add_text(s, f"{icon} {title}",
             left + Inches(0.4), top + Inches(0.1), w - Inches(0.5), Inches(0.42),
             font_size=Pt(14), font_color=DARK_NAVY, bold=True)
    add_text(s, subtitle,
             left + Inches(0.4), top + Inches(0.52), w - Inches(0.5), Inches(0.3),
             font_size=Pt(11), font_color=LIGHT_SLATE, italic=True)
    txBox = s.shapes.add_textbox(
        left + Inches(0.4), top + Inches(0.85), w - Inches(0.55), Inches(1.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, m in enumerate(mitigations):
        add_para(tf, f"✓  {m}", font_size=Pt(11), font_color=MID_SLATE,
                 space_before=Pt(3 if j > 0 else 0))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11  —  GO-TO-MARKET
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Go-to-Market Strategy", "Hub-first expansion: own a community, own a corridor.")
footer(s)

phases = [
    ("Phase 1\nQ3 2025", "Pilot Hubs",
     ["3 university campuses (target: 500 rides/month each)",
      "2 hospital networks (staff commute corridors)",
      "Recruit 50 Founding Drivers per hub",
      "NPS and safety incident baseline established"]),
    ("Phase 2\nQ1 2026", "City Corridors",
     ["Expand to commuter corridors in 3 metro areas",
      "Corporate hub B2B sales motion launched",
      "App Store featuring via safety & sustainability angle",
      "PR: 'The anti-Uber' positioning in mainstream media"]),
    ("Phase 3\nQ3 2026", "Platform Scale",
     ["API open to third-party hub operators",
      "International expansion: UK, Canada (regulatory-light)",
      "Kind Points marketplace with 10+ redemption partners",
      "Series A raise supported by unit economics data"]),
]

for i, (phase, title, bullets) in enumerate(phases):
    left = Inches(0.4 + i * 4.3)
    # Phase label
    add_rect(s, left, Inches(1.5), Inches(4.0), Inches(0.65), TEAL)
    add_text(s, phase, left, Inches(1.52), Inches(4.0), Inches(0.6),
             font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Arrow connector
    if i < 2:
        add_text(s, "→", left + Inches(4.05), Inches(1.65), Inches(0.3), Inches(0.4),
                 font_size=Pt(22), font_color=TEAL, bold=True)
    # Card
    bullet_card(s, title, bullets,
                left, Inches(2.25), Inches(4.0), Inches(4.75))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12  —  FINANCIALS & USE OF FUNDS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Financial Projections & Use of Funds", "Seed ask: $750K. Path to profitability by Month 24.")
footer(s)

# Projection table
proj_headers = ["", "Year 1", "Year 2", "Year 3"]
proj_rows = [
    ("Active Hubs",           "6",       "28",       "90"),
    ("Monthly Active Users",  "2,400",   "18,000",   "85,000"),
    ("Monthly Rides",         "8,000",   "65,000",   "340,000"),
    ("Gross Revenue",         "$96K",    "$1.04M",   "$6.1M"),
    ("Gross Margin",          "62%",     "68%",      "74%"),
    ("EBITDA",                "($480K)", "($120K)",  "$1.4M"),
]

col_w = Inches(2.4)
row_h = Inches(0.48)
tl = Inches(0.4)
tt = Inches(1.5)
for j, h in enumerate(proj_headers):
    add_rect(s, tl + j * col_w, tt, col_w, row_h, DARK_NAVY)
    add_text(s, h, tl + j * col_w + Inches(0.05), tt + Inches(0.08),
             col_w - Inches(0.1), Inches(0.32),
             font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(proj_rows):
    rt = tt + (i + 1) * row_h
    for j, cell in enumerate(row):
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        if j == 0: bg = RGBColor(0xE2, 0xE8, 0xF0)
        add_rect(s, tl + j * col_w, rt, col_w, row_h, bg,
                 line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(0.5))
        fc = ACCENT_RED if "(" in cell else (TEAL if j > 0 else DARK_NAVY)
        add_text(s, cell, tl + j * col_w + Inches(0.05), rt + Inches(0.08),
                 col_w - Inches(0.1), Inches(0.32),
                 font_size=Pt(12), font_color=fc,
                 bold=(j == 0), align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Use of funds
add_text(s, "Use of Seed Funds  ($750,000)",
         Inches(10.0), Inches(1.55), Inches(3.1), Inches(0.42),
         font_size=Pt(14), font_color=DARK_NAVY, bold=True)

fund_items = [
    ("Engineering & Product",  "35%", "$262K"),
    ("Growth & Hub Acquisition","25%", "$187K"),
    ("Safety & Legal",          "15%", "$112K"),
    ("Operations",              "15%", "$112K"),
    ("Reserve",                 "10%", "$75K"),
]
bar_top = Inches(2.1)
for i, (label, pct, amount) in enumerate(fund_items):
    y = bar_top + i * Inches(0.82)
    add_text(s, label, Inches(10.0), y, Inches(1.9), Inches(0.35),
             font_size=Pt(11.5), font_color=MID_SLATE)
    pct_val = int(pct.replace("%","")) / 100
    bar_total_w = Inches(2.4)
    add_rect(s, Inches(10.0), y + Inches(0.38), bar_total_w, Inches(0.25),
             RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, Inches(10.0), y + Inches(0.38), Inches(bar_total_w.inches * pct_val), Inches(0.25), TEAL)
    add_text(s, f"{pct}  {amount}", Inches(12.5), y + Inches(0.37),
             Inches(0.8), Inches(0.3),
             font_size=Pt(11), font_color=TEAL, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13  —  THE TEAM
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Team", "Mission-driven builders with skin in the game.")
footer(s)

# Founder card
add_rect(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(3.2), WHITE,
         line_color=TEAL, line_width=Pt(2))
add_rect(s, Inches(0.5), Inches(1.55), Inches(0.35), Inches(3.2), TEAL)
add_text(s, "Oluwafemi Adebayo",
         Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.55),
         font_size=Pt(20), font_color=DARK_NAVY, bold=True)
add_text(s, "Founder & CEO",
         Inches(1.0), Inches(2.28), Inches(4), Inches(0.35),
         font_size=Pt(13), font_color=TEAL, bold=True)
founder_points = [
    "Built the full KindRide platform end-to-end (mobile + backend + infra)",
    "Deep expertise: React Native, FastAPI, Supabase, AWS, AI/ML",
    "Identified the zero-detour gap from lived experience of unreliable transport",
    "Completed AWS Training certification programme",
]
txBox = s.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(5.1), Inches(1.9))
tf = txBox.text_frame
tf.word_wrap = True
for i, pt in enumerate(founder_points):
    add_para(tf, f"• {pt}", font_size=Pt(12), space_before=Pt(4 if i > 0 else 0))

# Advisory needs
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(3.2), WHITE,
         line_color=RGBColor(0xDB, 0xE4, 0xF5), line_width=Pt(1))
add_text(s, "Advisory & Hiring Plan",
         Inches(7.0), Inches(1.7), Inches(5.6), Inches(0.45),
         font_size=Pt(15), font_color=DARK_NAVY, bold=True)

hiring = [
    ("🎯", "Head of Growth",         "Hub BD + campus partnerships"),
    ("⚖️", "Transport Law Advisor",   "Regulatory strategy (US + UK)"),
    ("📱", "Senior Mobile Engineer",  "iOS optimisation + offline-first"),
    ("🤝", "Community Manager",       "Driver onboarding + retention"),
    ("💳", "Payments Lead",           "Stripe integration + financial compliance"),
]
for i, (icon, role, desc) in enumerate(hiring):
    y = Inches(2.25) + i * Inches(0.46)
    add_text(s, f"{icon}  {role}",
             Inches(7.0), y, Inches(2.5), Inches(0.38),
             font_size=Pt(12), font_color=DARK_NAVY, bold=True)
    add_text(s, desc,
             Inches(9.65), y, Inches(2.9), Inches(0.38),
             font_size=Pt(12), font_color=MID_SLATE)

add_text(s, "* Seed funding enables first 3 hires within 90 days of close.",
         Inches(6.8), Inches(4.85), Inches(6.0), Inches(0.35),
         font_size=Pt(10), font_color=LIGHT_SLATE, italic=True)

# Why founder story matters box
add_rect(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.85),
         RGBColor(0xCC, 0xFB, 0xF1))
add_rect(s, Inches(0.5), Inches(5.05), Inches(0.2), Inches(1.85), TEAL)
add_text(s, "Why a solo technical founder is a strength at this stage:",
         Inches(0.85), Inches(5.15), Inches(11.5), Inches(0.4),
         font_size=Pt(13), font_color=DARK_NAVY, bold=True)
add_text(s,
         "The entire product exists and functions in production. Zero agency fees, zero miscommunication, full IP ownership. "
         "The seed round funds growth and team — not product discovery. "
         "Every dollar invested goes to users, not building something that already exists.",
         Inches(0.85), Inches(5.6), Inches(12.0), Inches(1.1),
         font_size=Pt(12), font_color=MID_SLATE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14  —  THE ASK
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)

add_rect(s, 0, 0, Inches(0.5), SLIDE_H, TEAL)

add_text(s, "The Ask",
         Inches(1.2), Inches(0.7), Inches(11), Inches(0.9),
         font_size=Pt(48), font_color=WHITE, bold=True)

add_rect(s, Inches(1.2), Inches(1.65), Inches(10.5), Pt(3), TEAL)

add_text(s, "We are raising $750,000 in a Seed round.",
         Inches(1.2), Inches(1.8), Inches(11), Inches(0.6),
         font_size=Pt(22), font_color=TEAL, bold=True)

deliverables = [
    ("6 months", "Live pilots in 6 hubs. 2,400 MAU. 8,000 monthly rides. Safety record established."),
    ("12 months", "Series A metrics in sight: $1M ARR run-rate, 28 hubs, product-market fit proven."),
    ("18 months", "Profitable unit economics per hub. Expansion playbook documented. Team of 8."),
]
for i, (timeline, outcome) in enumerate(deliverables):
    y = Inches(2.55) + i * Inches(0.95)
    add_rect(s, Inches(1.2), y, Inches(1.5), Inches(0.75), TEAL)
    add_text(s, timeline, Inches(1.2), y + Inches(0.12), Inches(1.5), Inches(0.5),
             font_size=Pt(13), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(2.8), y, Inches(9.8), Inches(0.75),
             RGBColor(0x1E, 0x29, 0x3B))
    add_text(s, outcome, Inches(3.0), y + Inches(0.15), Inches(9.4), Inches(0.5),
             font_size=Pt(13), font_color=RGBColor(0xCB, 0xD5, 0xE1))

add_text(s, "KindRide is not just a better ride-hailing app.",
         Inches(1.2), Inches(5.55), Inches(11), Inches(0.5),
         font_size=Pt(18), font_color=WHITE, bold=True)
add_text(s,
         "It is a trust infrastructure for human mobility — "
         "the first platform to align the incentives of drivers, passengers, communities, and the planet.",
         Inches(1.2), Inches(6.05), Inches(11), Inches(0.65),
         font_size=Pt(14), font_color=LIGHT_SLATE)

add_text(s, "mcfemy@gmail.com  ·  kindride.app",
         Inches(1.2), Inches(6.9), Inches(6), Inches(0.35),
         font_size=Pt(12), font_color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out_path = "docs/KindRide_Investor_Pitch_Deck.pptx"
prs.save(out_path)
print(f"✓ Saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
