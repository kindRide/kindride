from pptx import Presentation
from pptx.util import Inches, Pt
import os

def add_slide(prs, layout_idx, title, content, notes=""):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    # Set Content
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = content[0] if content else ""
        for point in content[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    # Add Speaker Notes
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
        
    return slide

def generate_pitch_deck():
    prs = Presentation()
    
    # Slide 1: Title Slide (Layout 0)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "KindRide"
    slide.placeholders[1].text = "Ride-sharing built on community, trust, and zero detour.\n\nOluwafemi Adebayo, Founder & CEO | Seeking Seed Capital"
    slide.notes_slide.notes_text_frame.text = "Keep this brief. State your name, your company, and your one-sentence value proposition. 'Hi, I'm Femi, and KindRide is unlocking the 35% of daily urban vehicle capacity that currently travels empty, without charging passengers a dime.'"

    # Slide 2: The Problem
    add_slide(prs, 1, 
              "The current mobility model is broken", 
              ["Prohibitively Expensive: Uber/Lyft surge pricing locks low-income and daily commuters out of reliable transport.",
               "Environmental Waste: Ride-hail deadheading means 40% of miles are driven empty. 35% of daily urban vehicle capacity travels completely empty.",
               "Broken Trust: Existing platforms commoditize human connection. Safety incidents have eroded confidence."],
              "Don't spend too long here. VCs know Uber is expensive. Focus on the '35% of urban capacity travels empty' stat. That is the wasted resource you are going to capture.")

    # Slide 3: The Solution
    add_slide(prs, 1, 
              "Revealing trips that are already happening", 
              ["Zero-Detour Algorithm: We match passengers with drivers ALREADY traveling in their exact direction.",
               "Zero Fare: Because the marginal cost to the driver is zero, the ride is free.",
               "The Currency is Kindness: Drivers are compensated in 'Kind Points' (social capital) and voluntary tips, rather than extracted wages."],
              "This is your 'Aha!' moment. Emphasize that KindRide doesn't create trips; it reveals them. The driver was going there anyway.")

    # Slide 4: The Magic (How It Works)
    add_slide(prs, 1, 
              "A frictionless, premium experience", 
              ["Smart Matching: Ranks drivers by Proximity, Heading Alignment, and Urgency.",
               "Multi-Leg Handoffs: If no single driver is going the full distance, the system intelligently chains trips.",
               "UX Polish: Shimmer loading, haptic feedback, 'vibe' selectors (silent vs. chatty), and live GPS tracking."],
              "VCs love to see that a solo founder understands UX. Mention that the app feels like a $100M company built it. Call out the Multi-Leg feature—it proves you are solving complex routing problems, not just cloning Uber.")

    # Slide 5: The Safety & Trust Architecture
    add_slide(prs, 1, 
              "Safety isn't an afterthought; it's our infrastructure", 
              ["Verified Identity: Mandatory Stripe KYC Identity verification for drivers.",
               "In-App Trip Recording: Audio/telemetry is recorded locally and auto-deleted after 72 hours unless flagged.",
               "Live SOS & Tracking: One-tap emergency alerts and shareable live-tracking links.",
               "Zero-Knowledge Route Commitments: Cryptographic proof that drivers stayed on the route."],
              "Investors worry about liability. This slide proves you are mitigating risk better than the incumbents. Mention the 72-hour ephemeral recording—it balances safety perfectly with privacy.")

    # Slide 6: Gamifying Generosity
    add_slide(prs, 1, 
              "The 'Kind Points' Economy", 
              ["Earning: Drivers earn a base of 10 points + 1 point per mile + streaks + 5-star bonuses.",
               "Real Impact: Translating metrics to impact: '150 points = ~15 people helped.'",
               "Redemptions: Points can be donated to local shelters, unlock profile badges, or be used to 'Pay It Forward.'"],
              "Explain that humans are wired for status and altruism. Points act as a 'social credit score' that drivers can be proud of in their communities.")

    # Slide 7: Go-To-Market
    add_slide(prs, 1, 
              "The 'Walled Garden' Strategy: Community Hubs", 
              ["The Problem: City-wide launches fail due to low liquidity (chicken-and-egg).",
               "Our Solution: We launch inside B2B Community Hubs (Universities, Hospitals, Corporate Campuses).",
               "How it works: Users auto-join via their .edu or .org email. This creates dense, high-trust micro-networks with captive supply and demand."],
              "Deliver this with extreme confidence. This is the most important slide for a VC. Launching city-wide is too expensive. Launching at a specific university requires almost zero marketing spend because the community is already dense.")

    # Slide 8: The Business Model
    add_slide(prs, 1, 
              "Multiple Monetization Levers", 
              ["B2B Hub Subscriptions (Primary): Institutions pay a monthly SaaS fee to maintain private, verified commuter networks.",
               "Voluntary Tipping & Premium: Stripe-powered passenger tipping (we take standard processing margins). Optional 'Safety Premium' tiers.",
               "Kind Points Marketplace & Data: Local businesses sponsor Point redemptions; anonymized mobility data sold to urban planners."],
              "Make it clear: 'We do not take a 40% cut of a driver's wage like Uber. We monetize the infrastructure and the community access.'")

    # Slide 9: Market Opportunity
    add_slide(prs, 1, 
              "A massive market, ripe for a new paradigm", 
              ["TAM (Total Addressable Market): $285B Global Ride-Hail Market.",
               "SAM (Serviceable Market): $42B Community & Commuter ride-sharing (US/UK).",
               "SOM (Serviceable Obtainable Market): $1.2B (Year 5 capture via Campus + Metro Corridors)."],
              "Don't over-explain the numbers. Just show that the market is massive, and capturing even a tiny fraction of daily commuters yields a billion-dollar company.")

    # Slide 10: Competitive Landscape
    add_slide(prs, 1, 
              "Where KindRide Wins", 
              ["Waze Carpool: Rigid and scheduled (Discontinued).",
               "Uber/Lyft: Transactional, expensive, 40% driver take-rate.",
               "KindRide: Ad-hoc, real-time, multi-leg, community-verified transit.",
               "Key Moats: Zero-Detour Algorithm, Community Hub Walled Gardens, Zero Cost to Passenger."],
              "'Waze tried carpooling, but it was rigid and scheduled. Uber is transactional and expensive. We are the only platform offering ad-hoc, real-time, multi-leg, community-verified transit.'")

    # Slide 11: The Team
    add_slide(prs, 1, 
              "Built by a solo technical powerhouse", 
              ["Oluwafemi Adebayo – Founder & CEO",
               "Built the entire KindRide platform end-to-end (React Native, Python/FastAPI, Supabase, PostGIS).",
               "AWS Certified, deep expertise in scalable cloud architecture.",
               "The Founder Advantage: Zero agency fees, zero technical debt, 100% IP ownership."],
              "Brag about yourself here. VCs invest in founders first. Tell them: 'Because I built this myself, I have an execution speed that a team of 5 engineers at a standard startup couldn't match.'")

    # Slide 12: The Ask & Use of Funds
    add_slide(prs, 1, 
              "Raising $750,000 Seed", 
              ["Runway: 18–24 months.",
               "Goals for this round:",
               "  • Launch 6 Pilot Hubs (Universities/Hospitals).",
               "  • Reach 2,400 MAU and 8,000 monthly rides.",
               "  • Achieve $1M ARR run-rate to unlock Series A.",
               "Use of Funds: 35% Eng/Product, 25% Hub Acquisition, 15% Safety/Legal, 15% Ops, 10% Reserve."],
              "End with a strong closing statement. 'KindRide isn't just a better ride-hailing app. It is a trust infrastructure for human mobility. I'd love to partner with you to build it.'")

    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), "KindRide_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    generate_pitch_deck()